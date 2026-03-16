# SPDX-FileCopyrightText: 2024 MarkDowner Team
#
# SPDX-License-Identifier: MIT

"""Fixture corpus smoke tests."""

import importlib.util
import io
from pathlib import Path
import zipfile

import pytest

from markdowner import MarkDowner
from markdowner._stream_info import StreamInfo
from markdowner.converters._docx_converter import DocxConverter
from markdowner.converters._pdf_converter import PdfConverter
from markdowner.converters._xlsx_converter import XlsxConverter


FIXTURES = Path(__file__).parent / "test_files"


def _has_module(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def _build_docx_bytes(text: str) -> bytes:
    archive = io.BytesIO()
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body><w:p><w:r><w:t>"
        f"{text}"
        "</w:t></w:r></w:p></w:body></w:document>"
    )
    with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            "</Types>",
        )
        zf.writestr(
            "_rels/.rels",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
            'Target="word/document.xml"/>'
            "</Relationships>",
        )
        zf.writestr("word/document.xml", document_xml)
    return archive.getvalue()


def _build_pdf_bytes(text: str) -> bytes:
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 144] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        (
            b"<< /Length %d >>\nstream\n%s\nendstream"
            % (
                len(f"BT\n/F1 18 Tf\n72 100 Td\n({text}) Tj\nET".encode("latin-1")),
                f"BT\n/F1 18 Tf\n72 100 Td\n({text}) Tj\nET".encode("latin-1"),
            )
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    chunks = [b"%PDF-1.4\n"]
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(sum(len(chunk) for chunk in chunks))
        chunks.append(f"{index} 0 obj\n".encode("ascii"))
        chunks.append(obj)
        chunks.append(b"\nendobj\n")

    xref_offset = sum(len(chunk) for chunk in chunks)
    chunks.append(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    chunks.append(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        chunks.append(f"{offset:010d} 00000 n \n".encode("ascii"))
    chunks.append(
        f"trailer\n<< /Root 1 0 R /Size {len(objects) + 1} >>\nstartxref\n{xref_offset}\n%%EOF\n".encode(
            "ascii"
        )
    )
    return b"".join(chunks)


def test_fixture_directory_is_populated():
    names = sorted(path.name for path in FIXTURES.iterdir())
    assert names == [
        "rtf",
        "rtf_csv_like_sample.rtf",
        "rtf_expected",
        "sample-attachment.msg",
        "sample.csv",
        "sample.docx",
        "sample.html",
        "sample.msg",
        "sample.pdf",
        "sample.rtf",
        "sample.txt",
        "sample.xls",
        "sample.xlsx",
        "sample.zip",
    ]


def test_text_html_csv_and_zip_fixtures_convert():
    md = MarkDowner()

    text_result = md.convert(FIXTURES / "sample.txt")
    html_result = md.convert(FIXTURES / "sample.html")
    csv_result = md.convert(FIXTURES / "sample.csv")
    zip_result = md.convert(FIXTURES / "sample.zip")

    assert "Fixture text line." in text_result.text_content
    assert "# Fixture Title" in html_result.text_content
    assert "Alice" in csv_result.text_content
    assert "zip fixture text" in zip_result.text_content


@pytest.mark.skipif(not _has_module("mammoth"), reason="mammoth dependency not installed")
def test_docx_fixture_converts_end_to_end():
    converter = DocxConverter()
    payload = _build_docx_bytes("Fixture DOCX text")

    result = converter.convert(io.BytesIO(payload), StreamInfo(extension=".docx"))

    assert "Fixture DOCX text" in result.text_content


@pytest.mark.skipif(
    not (_has_module("pandas") and _has_module("openpyxl")),
    reason="pandas/openpyxl dependencies not installed",
)
def test_xlsx_fixture_converts_end_to_end(tmp_path):
    from openpyxl import Workbook

    fixture = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Fixture Sheet"
    sheet.append(["name", "value"])
    sheet.append(["alpha", 1])
    workbook.save(fixture)

    with fixture.open("rb") as fh:
        result = XlsxConverter().convert(
            fh,
            StreamInfo(
                extension=".xlsx",
                mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ),
        )

    assert "Fixture Sheet" in result.text_content
    assert "alpha" in result.text_content


@pytest.mark.skipif(not _has_module("pdfminer"), reason="pdfminer.six dependency not installed")
def test_pdf_fixture_converts_end_to_end():
    converter = PdfConverter()
    payload = _build_pdf_bytes("Fixture PDF line one.")
    result = converter.convert(io.BytesIO(payload), StreamInfo(extension=".pdf"))

    assert "Fixture PDF line one." in result.text_content


@pytest.mark.skipif(not _has_module("pptx"), reason="python-pptx dependency not installed")
def test_generated_pptx_fixture_converts_end_to_end(tmp_path):
    from pptx import Presentation

    fixture = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Fixture PPTX title"
    slide.placeholders[1].text = "Fixture PPTX body"
    presentation.save(fixture)

    result = MarkDowner().convert(fixture)

    assert "Fixture PPTX title" in result.text_content
    assert "Fixture PPTX body" in result.text_content


@pytest.mark.skipif(not _has_module("ebooklib"), reason="ebooklib dependency not installed")
def test_generated_epub_fixture_converts_end_to_end(tmp_path):
    from ebooklib import epub

    fixture = tmp_path / "sample.epub"
    book = epub.EpubBook()
    book.set_identifier("fixture-book")
    book.set_title("Fixture EPUB title")
    chapter = epub.EpubHtml(title="Intro", file_name="intro.xhtml", content="<h1>Intro</h1><p>Fixture EPUB body</p>")
    book.add_item(chapter)
    book.toc = (chapter,)
    book.spine = ["nav", chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(fixture), book)

    result = MarkDowner().convert(fixture)

    assert "# Fixture EPUB title" in result.text_content
    assert "Fixture EPUB body" in result.text_content


@pytest.mark.skipif(
    not (FIXTURES / "sample.xls").exists(),
    reason="sample.xls fixture is not present",
)
def test_xls_fixture_converts_end_to_end():
    result = MarkDowner().convert(FIXTURES / "sample.xls")
    assert result.text_content.strip()


@pytest.mark.skipif(
    not (FIXTURES / "sample.msg").exists(),
    reason="sample.msg fixture is not present in-repo and MSG generation is not available in the test stack",
)
def test_msg_fixture_converts_end_to_end():
    result = MarkDowner().convert(FIXTURES / "sample.msg")
    assert result.text_content.strip()


@pytest.mark.skipif(
    not (FIXTURES / "sample-attachment.msg").exists(),
    reason="sample-attachment.msg fixture is not present",
)
def test_msg_attachment_fixture_converts_end_to_end():
    """Test that MSG with attachment converts main body and handles attachment."""
    result = MarkDowner().convert(FIXTURES / "sample-attachment.msg")
    # Main body should be non-empty
    assert result.text_content.strip()
    # Should have attachment metadata (even if conversion fails, main body should work)
    # The conversion should not raise, and attachments should be in metadata
    # If attachment conversion fails, it should be non-fatal
    assert "attachment_outputs" in result.metadata or "attachment skipped" in str(result.warnings)
