# SPDX-FileCopyrightText: 2024 MarkDowner Team
#
# SPDX-License-Identifier: MIT

"""PPTX converter."""

from typing import BinaryIO

from .._base_converter import DocumentConverter, DocumentConverterResult
from .._limits import DEFAULT_LIMITS, Limits
from .._sandbox import ParserSandboxLimits, run_in_subprocess
from .._stream_info import StreamInfo
from .._temp_utils import materialize_stream_to_temp_path
from .._exceptions import MissingDependencyException
from ._zip_package_helpers import zip_has_members


def _convert_pptx_to_markdown(pptx_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        slides.append(f"## Slide {i}\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slides.append(shape.text)
                slides.append("\n")

    return "\n".join(slides)


class PptxConverter(DocumentConverter):
    """Converter for PPTX files."""

    def __init__(self, limits: Limits = DEFAULT_LIMITS):
        self._limits = limits

    def accepts(
        self,
        stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs
    ) -> bool:
        """Check if this is a PPTX file."""
        mimetype = stream_info.mimetype or ""
        extension = (stream_info.extension or "").lower()

        if "presentation" in mimetype or "powerpoint" in mimetype:
            return True

        if extension != ".pptx" and mimetype:
            return False

        return zip_has_members(
            stream,
            "[Content_Types].xml",
            "ppt/presentation.xml",
            max_entries=self._limits.max_zip_metadata_entries,
            max_metadata_bytes=self._limits.max_zip_metadata_scan_bytes,
        )

    def convert(
        self,
        stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs
    ) -> DocumentConverterResult:
        """Convert PPTX to Markdown."""
        try:
            from pptx import Presentation
        except ImportError:
            raise MissingDependencyException(
                "python-pptx is required for PPTX conversion. Install with: pip install markdowner[pptx]"
            )

        del Presentation

        try:
            with materialize_stream_to_temp_path(stream, ".pptx") as tmp_path:
                text = run_in_subprocess(
                    _convert_pptx_to_markdown,
                    str(tmp_path),
                    limits=ParserSandboxLimits(),
                )
        except Exception as e:
            raise RuntimeError(f"PPTX conversion failed: {e}") from e

        return DocumentConverterResult(text_content=text)
